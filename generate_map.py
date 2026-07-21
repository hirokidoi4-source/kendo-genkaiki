from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# プレゼンテーション作成 (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# タイトルテキスト
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "けやき祭 前庭会場案内図（編集用）"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(30, 41, 59)

# 説明テキスト
txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.333), Inches(0.5))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "※各枠・テキスト・図形は自由に変更・複製・移動が可能です。"
p2.font.size = Pt(14)
p2.font.color.rgb = RGBColor(100, 116, 139)

# サンプルブース枠（自由に変更・複製可能）
def add_booth(slide, left, top, width, height, title, detail, bg_rgb, border_rgb):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_rgb
    shape.line.color.rgb = border_rgb
    shape.line.width = Pt(2)
    
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = detail
    p2.font.size = Pt(10)
    p2.font.color.rgb = RGBColor(51, 65, 85)
    p2.alignment = PP_ALIGN.CENTER

# テント枠の配置サンプル
add_booth(slide, 1.0, 2.0, 1.8, 1.0, "【3-1】", "店舗名：未定\n出し物：未定", RGBColor(255, 255, 255), RGBColor(37, 99, 235))
add_booth(slide, 3.0, 2.0, 1.8, 1.0, "【3-2】", "店舗名：未定\n出し物：未定", RGBColor(255, 255, 255), RGBColor(37, 99, 235))
add_booth(slide, 5.0, 2.0, 2.2, 1.2, "【飲食スペース】", "自由席テーブル\n配置変更自由", RGBColor(254, 243, 199), RGBColor(217, 119, 6))
add_booth(slide, 7.5, 2.0, 1.8, 1.0, "★ 入場口", "受付・記入台前\n← 進入順路", RGBColor(252, 165, 165), RGBColor(220, 38, 38))

# ファイル保存
prs.save("けやき祭_会場案内図_編集ベース.pptx")
print("PowerPointファイルを生成しました！")
